"""
Serviço para converter documentos (PDF/Excel/DOCX/PPTX/Imagens) para Markdown.
MVP 100% FREE: Usa apenas soluções leves e gratuitas.

- PDF (texto nativo): PyPDF2/pdfplumber - Extrai texto de PDFs com texto nativo
- Excel: pandas/openpyxl - Conversão completa de planilhas
- DOCX: python-docx - Conversão preservando estrutura
- PPTX: python-pptx - Conversão com estrutura de slides
- PDF (escaneado/OCR): Suporte básico via OCR Service (opcional)
- Imagens (JPEG, PNG, GIF, BMP, TIFF, WEBP): OCR via pytesseract (opcional)

LIMITAÇÕES:
- PDFs escaneados requerem OCR (biblioteca opcional)
- PDFs com imagens não extraem texto das imagens automaticamente
- Imagens requerem OCR (biblioteca opcional)
"""
from typing import Optional
from pathlib import Path
from io import BytesIO
from loguru import logger
import re

from app.services.pdf_converter import convert_pdf_to_markdown
from app.services.excel_converter import convert_excel_to_markdown

# Importar OCR Service (opcional)
try:
    from app.services.ocr_service import OCRService
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    OCRService = None

# Importar Multimodal Service (Vision Advanced)
try:
    from src.features.vision.multimodal_service import multimodal_service
    VISION_AVAILABLE = True
except ImportError:
    VISION_AVAILABLE = False
    multimodal_service = None

# Verificar disponibilidade de bibliotecas
try:
    import PyPDF2
    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False
    logger.warning("PyPDF2 não instalado. PDF não será suportado.")

try:
    EXCEL_AVAILABLE = True
    import pandas as pd
except ImportError:
    EXCEL_AVAILABLE = False
    logger.warning("pandas/openpyxl não instalados. Excel não será suportado.")

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx não instalado. DOCX não será suportado.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx não instalado. PPTX não será suportado.")


class DocumentConverterService:
    """Serviço para converter documentos para Markdown (MVP 100% FREE)."""
    
    def __init__(self, enable_ocr: bool = False):
        """
        Inicializa o conversor de documentos.
        
        Args:
            enable_ocr: Se True, tenta usar OCR para PDFs escaneados (requer bibliotecas)
        """
        # Inicializar OCR Service se solicitado e disponível
        self.enable_ocr = enable_ocr and OCR_AVAILABLE
        self.ocr_service = None
        if self.enable_ocr:
            try:
                self.ocr_service = OCRService()
                if self.ocr_service.is_ocr_available():
                    logger.info("✅ OCR habilitado para PDFs escaneados e imagens")
                else:
                    logger.warning("OCR solicitado mas não disponível (bibliotecas faltando)")
                    self.enable_ocr = False
            except Exception as e:
                logger.warning(f"Erro ao inicializar OCR: {e}")
                self.enable_ocr = False

        supported_formats = []
        
        if PYPDF2_AVAILABLE:
            supported_formats.append("PDF")
        if EXCEL_AVAILABLE:
            supported_formats.append("Excel")
        if DOCX_AVAILABLE:
            supported_formats.append("DOCX")
        if PPTX_AVAILABLE:
            supported_formats.append("PPTX")
        if self.enable_ocr and self.ocr_service and self.ocr_service.is_ocr_available():
            supported_formats.append("Imagens (OCR)")
        if VISION_AVAILABLE:
            supported_formats.append("Visão Multimodal (Gemini 1.5 Pro)")
        
        if supported_formats:
            logger.info(f"✅ DocumentConverterService inicializado - Formatos suportados: {', '.join(supported_formats)}")
        else:
            logger.warning("DocumentConverterService inicializado - Nenhum formato suportado (instale dependências)")
        
        # Logs de avisos para formatos não disponíveis
        if not PYPDF2_AVAILABLE:
            logger.warning("PDF não disponível - PyPDF2 não instalado")
        if not EXCEL_AVAILABLE:
            logger.warning("Excel não disponível - pandas/openpyxl não instalados")
        if not DOCX_AVAILABLE:
            logger.warning("DOCX não disponível - python-docx não instalado")
        if not PPTX_AVAILABLE:
            logger.warning("PPTX não disponível - python-pptx não instalado")
    
    async def convert_file(self, file_path: str) -> Optional[str]:
        """
        Converte arquivo do sistema de arquivos para Markdown.
        
        Args:
            file_path: Caminho do arquivo a converter
            
        Returns:
            str: Conteúdo Markdown ou None se erro
        """
        path = Path(file_path)
        
        if not path.exists():
            logger.error(f"Arquivo não encontrado: {file_path}")
            return None
        
        suffix = path.suffix.lower()
        
        # Verificar se é Excel
        if suffix in ['.xlsx', '.xls']:
            if not EXCEL_AVAILABLE:
                logger.error("Excel não suportado - pandas/openpyxl não instalados")
                return None
            return convert_excel_to_markdown(path.read_bytes(), path.name)
        
        # Verificar se é PDF
        if suffix == '.pdf':
            if not PYPDF2_AVAILABLE:
                logger.error("PDF não suportado - PyPDF2 não instalado")
                return None
            return convert_pdf_to_markdown(path.read_bytes(), path.name)
        
        # Formatos não suportados
        logger.warning(f"Formato {suffix} não suportado no MVP (apenas PDF texto nativo e Excel)")
        return None
    
    async def convert_bytes(self, file_content: bytes, filename: str) -> Optional[str]:
        """
        Converte arquivo a partir de bytes (para upload) para Markdown.
        
        Args:
            file_content: Conteúdo do arquivo em bytes
            filename: Nome do arquivo (para detectar formato)
            
        Returns:
            str: Conteúdo Markdown ou None se erro
        """
        suffix = Path(filename).suffix.lower()
        
        # Roteamento: Excel usa conversão manual
        if suffix in ['.xlsx', '.xls']:
            if not EXCEL_AVAILABLE:
                logger.error("Excel não suportado - pandas/openpyxl não instalados")
                return None
            return convert_excel_to_markdown(file_content, filename)
        
        # PDF
        if suffix == '.pdf':
            if not PYPDF2_AVAILABLE:
                logger.error("PDF não suportado - PyPDF2 não instalado")
                return None
            
            # Tentar conversão normal primeiro
            markdown_content = convert_pdf_to_markdown(file_content, filename)
            
            # Se falhou e OCR está habilitado, tentar OCR
            if not markdown_content and self.enable_ocr and self.ocr_service:
                logger.info(f"Conversão PDF normal falhou, tentando OCR para: {filename}")
                markdown_content = self.ocr_service.process_scanned_pdf(file_content, filename)
            
            return markdown_content
        
        # DOCX
        if suffix == '.docx':
            if not DOCX_AVAILABLE:
                logger.error("DOCX não suportado - python-docx não instalado")
                return None
            return self.convert_docx_to_markdown(file_content, filename)
        
        # PPTX
        if suffix == '.pptx':
            if not PPTX_AVAILABLE:
                logger.error("PPTX não suportado - python-pptx não instalado")
                return None
            return self.convert_pptx_to_markdown(file_content, filename)
        
        # Imagens (JPEG, PNG, GIF, BMP, TIFF, WEBP)
        if suffix in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp']:
            # Prioridade 1: Visão Multimodal (Gemini)
            if VISION_AVAILABLE and multimodal_service:
                return await self.convert_image_via_vision(file_content, filename)
            
            # Prioridade 2: OCR padrão
            if not self.enable_ocr or not self.ocr_service:
                logger.error("Imagens não suportadas - OCR não habilitado e Visão não disponível")
                return None
            return self.convert_image_to_markdown(file_content, filename)
        
        # Formatos não suportados
        logger.warning(f"Formato {suffix} não suportado")
        return None
    
    def convert_docx_to_markdown(self, file_content: bytes, filename: str) -> Optional[str]:
        """
        Converte DOCX para Markdown preservando estrutura básica.
        
        Args:
            file_content: Conteúdo do arquivo em bytes
            filename: Nome do arquivo (para logging)
            
        Returns:
            str: Conteúdo Markdown ou None se erro
        """
        if not DOCX_AVAILABLE:
            logger.error("DOCX não suportado - python-docx não instalado")
            return None
        
        try:
            doc = Document(BytesIO(file_content))
            markdown_lines = []
            
            # Processar parágrafos
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # Detectar cabeçalhos baseado no estilo
                style_name = para.style.name.lower()
                if 'heading' in style_name or 'title' in style_name:
                    # Extrair nível do cabeçalho (Heading 1, Heading 2, etc.)
                    level_match = re.search(r'heading\s*(\d+)', style_name)
                    if level_match:
                        level = int(level_match.group(1))
                    elif 'title' in style_name:
                        level = 1
                    else:
                        level = 2
                    
                    markdown_lines.append(f"{'#' * level} {text}\n")
                else:
                    markdown_lines.append(f"{text}\n")
            
            # Processar tabelas
            for table in doc.tables:
                markdown_lines.append(self._table_to_markdown(table))
                markdown_lines.append("\n")
            
            result = "\n".join(markdown_lines).strip()
            
            if not result:
                logger.warning(f"DOCX '{filename}' convertido mas sem conteúdo extraído")
                return None
            
            logger.info(f"DOCX '{filename}' convertido: {len(result)} caracteres")
            return result
            
        except Exception as e:
            logger.error(f"Erro ao converter DOCX '{filename}': {e}")
            import traceback
            logger.error(traceback.format_exc())
            return None
    
    def convert_pptx_to_markdown(self, file_content: bytes, filename: str) -> Optional[str]:
        """
        Converte PPTX para Markdown com estrutura de slides.
        
        Args:
            file_content: Conteúdo do arquivo em bytes
            filename: Nome do arquivo (para logging)
            
        Returns:
            str: Conteúdo Markdown ou None se erro
        """
        if not PPTX_AVAILABLE:
            logger.error("PPTX não suportado - python-pptx não instalado")
            return None
        
        try:
            prs = Presentation(BytesIO(file_content))
            markdown_lines = []
            
            for i, slide in enumerate(prs.slides, 1):
                markdown_lines.append(f"# Slide {i}\n")
                
                # Processar formas na slide
                for shape in slide.shapes:
                    if not hasattr(shape, "text"):
                        continue
                    
                    text = shape.text.strip()
                    if not text:
                        continue
                    
                    # Detectar títulos (geralmente são as primeiras formas ou têm formato específico)
                    # Simplificação: primeira forma não vazia é título
                    if i == 1 and len(markdown_lines) == 1:
                        # Primeira forma da primeira slide pode ser título
                        markdown_lines.append(f"## {text}\n")
                    elif text.upper() in ["TÍTULO", "TITLE", "TITULO"] or len(text) < 50:
                        # Textos curtos ou que contêm "título" podem ser títulos
                        markdown_lines.append(f"## {text}\n")
                    else:
                        markdown_lines.append(f"{text}\n")
                
                markdown_lines.append("\n")
            
            result = "\n".join(markdown_lines).strip()
            
            if not result:
                logger.warning(f"PPTX '{filename}' convertido mas sem conteúdo extraído")
                return None
            
            logger.info(f"PPTX '{filename}' convertido: {len(result)} caracteres, {len(prs.slides)} slides")
            return result
            
        except Exception as e:
            logger.error(f"Erro ao converter PPTX '{filename}': {e}")
            import traceback
            logger.error(traceback.format_exc())
            return None
    
    def convert_image_to_markdown(self, image_bytes: bytes, filename: str) -> Optional[str]:
        """
        Converte imagem para Markdown usando OCR.
        
        Args:
            image_bytes: Conteúdo da imagem em bytes
            filename: Nome do arquivo (para logging)
            
        Returns:
            str: Conteúdo Markdown ou None se erro
        """
        if not self.enable_ocr or not self.ocr_service:
            logger.error("OCR não habilitado ou não disponível para processar imagens")
            return None
        
        try:
            markdown_content = self.ocr_service.process_image(image_bytes, filename)
            
            if not markdown_content:
                logger.warning(f"Imagem '{filename}' processada mas sem texto extraído")
                return None
            
            logger.info(f"Imagem '{filename}' convertida: {len(markdown_content)} caracteres")
            return markdown_content
            
        except Exception as e:
            logger.error(f"Erro ao converter imagem '{filename}': {e}")
            import traceback
            logger.error(traceback.format_exc())
            return None

    async def convert_image_via_vision(self, image_bytes: bytes, filename: str) -> Optional[str]:
        """
        Converte imagem para Markdown usando Visão Multimodal Avançada (Gemini).
        """
        if not VISION_AVAILABLE or not multimodal_service:
            return None

        try:
            logger.info(f"Usando Gemini Vision para processar: {filename}")
            
            # Obter descrição e análise estruturada
            description = await multimodal_service.describe_image(image_bytes)
            analysis = await multimodal_service.analyze_document_page(image_bytes)

            # Montar Markdown rico
            md = [
                f"## Análise Visual de {filename}",
                f"**Descrição:** {description}\n",
            ]

            if analysis:
                if "summary" in analysis:
                    md.append(f"### Resumo Analítico\n{analysis['summary']}\n")
                if "tables" in analysis and analysis["tables"]:
                    md.append("### Dados Extraídos (Tabelas)")
                    for i, table in enumerate(analysis["tables"], 1):
                        md.append(f"#### Tabela {i}")
                        if isinstance(table, list) and len(table) > 0:
                            # Tabela pode ser lista de listas ou lista de dicts
                            if isinstance(table[0], dict):
                                # Lista de dicionários
                                all_keys = []
                                for row in table:
                                    if isinstance(row, dict):
                                        for k in row.keys():
                                            if k not in all_keys:
                                                all_keys.append(k)
                                
                                if all_keys:
                                    md.append("| " + " | ".join(all_keys) + " |")
                                    md.append("| " + " | ".join(["---"] * len(all_keys)) + " |")
                                    for row in table:
                                        if isinstance(row, dict):
                                            vals = [str(row.get(k, "")) for k in all_keys]
                                            md.append("| " + " | ".join(vals) + " |")
                            elif isinstance(table[0], list):
                                # Lista de listas (row-based)
                                for row_idx, row in enumerate(table):
                                    md.append("| " + " | ".join(str(cell) for cell in row) + " |")
                                    if row_idx == 0: # Header
                                        md.append("| " + " | ".join(["---"] * len(row)) + " |")
                        else:
                            md.append(str(table))
                        md.append("")
                if "charts" in analysis and analysis["charts"]:
                    md.append(f"### Insights de Gráficos\n{str(analysis['charts'])}\n")
                if "alerts" in analysis and analysis["alerts"]:
                    md.append(f"### ⚠️ Alertas e Pontos Críticos")
                    if isinstance(analysis["alerts"], list):
                        for alert in analysis["alerts"]:
                            md.append(f"- {alert}")
                    else:
                        md.append(str(analysis["alerts"]))
                    md.append("")

            return "\n".join(md)

        except Exception as e:
            logger.error(f"Erro ao processar imagem via Vision: {e}")
            # Fallback para OCR tradicional se disponível
            if self.enable_ocr and self.ocr_service:
                logger.info("Tentando fallback para OCR tradicional...")
                return self.convert_image_to_markdown(image_bytes, filename)
            return None
    
    def _table_to_markdown(self, table) -> str:
        """
        Converte tabela do Word para Markdown.
        
        Args:
            table: Tabela do python-docx
            
        Returns:
            str: Tabela em formato Markdown
        """
        try:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            
            if not rows:
                return ""
            
            # Criar separador
            num_cols = len(table.rows[0].cells)
            separator = "|" + "|".join(["---"] * num_cols) + "|"
            
            # Primeira linha é cabeçalho, resto são dados
            if len(rows) > 1:
                return "\n".join([rows[0], separator] + rows[1:])
            else:
                return rows[0]
                
        except Exception as e:
            logger.warning(f"Erro ao converter tabela para Markdown: {e}")
            return ""
